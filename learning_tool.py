import os
import PyPDF2
from pptx import Presentation
import requests
import spacy

# Set Mistral API Key
MISTRAL_API_KEY = "3qE8fjR3Kg58syTtw2S9vk7Rw9Y6okch"
if not MISTRAL_API_KEY:
    raise ValueError("Mistral API key is missing. Set it as an environment variable.")

MISTRAL_API_URL = "https://api.mistral.ai/v1/chat/completions"

# Chunk size limit
CHUNK_SIZE = 4000  # Adjust as needed based on API token limits

# Extract text from PDF
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

# Extract text from PowerPoint
def extract_text_from_ppt(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

# Break text into chunks
def chunk_text(text, chunk_size=CHUNK_SIZE):
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

# Call Mistral AI API
def query_mistral(prompt):
    headers = {
        "Authorization": f"Bearer {MISTRAL_API_KEY}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": "mistral-small",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7,
    }
    response = requests.post(MISTRAL_API_URL, headers=headers, json=payload)
    if response.status_code == 200:
        return response.json()["choices"][0]["message"]["content"]
    else:
        return f"Error: {response.status_code}, {response.text}"

# Summarize text using Mistral AI
def summarize_text(text):
    chunks = chunk_text(text)
    summaries = []
    for chunk in chunks:
        prompt = f"Extract and list all key topics covered in the following text as a single paragraph, separated by commas. Exclude any headings, irrelevant details, or extra information:\n{chunk}"
        summaries.append(query_mistral(prompt))
    print(" ".join(summaries))
    return " ".join(summaries)

# Extract key concepts using spaCy
def extract_key_concepts(text):
    nlp = spacy.load("en_core_web_sm")
    doc = nlp(text)
    concepts = set()
    for ent in doc.ents:
        concepts.add((ent.text, ent.label_))
    return concepts

# Generate questions using Mistral AI
def generate_questions(text,keywords, question_type):
    prompt_map = {
        "one_word": (
            "Generate 15 fact-based questions that can be answered in a single word. "
            "Ensure that the questions test key concepts, definitions, and important terms from the following text:\n"
        ),
        "true_false": (
            "Generate 15 well-structured true or false questions that challenge the reader's understanding of key facts, "
            "concepts, and statements from the following text. Avoid ambiguous or overly obvious statements:\n"
        ),
        "3_mark": (
            "Generate 15 thought-provoking 3-mark questions that test comprehension, analysis, and explanation of key ideas "
            "from the following text. Ensure that the questions encourage detailed yet concise responses:\n"
        ),
        "5_mark": (
            "Generate 15 in-depth, higher-order thinking questions that require critical analysis and a comprehensive response. "
            "The questions should encourage discussion, explanation of concepts, and detailed understanding of the following text:\n"
        )
    }

    if question_type not in prompt_map:
        return "Invalid question type!"
    
    chunks = chunk_text(text)
    questions = []
    for chunk in chunks:
        prompt = prompt_map[question_type] + chunk
        questions.append(query_mistral(prompt))
        
    for concept, label in keywords:
        prompt = prompt_map[question_type] + concept
        questions.append(query_mistral(prompt))
    
    return "\n".join(questions)

# Main function
def main():
    file_path = input("Enter the file path (PDF or PPT): ").strip()
    if file_path.endswith(".pdf"):
        text = extract_text_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        text = extract_text_from_ppt(file_path)
    else:
        print("Unsupported file format!")
        return

    summary = summarize_text(text)
    print("\n=== Summary ===")
    print(summary)

    concepts = extract_key_concepts(summary)
    print("\n=== Key Concepts ===")
    for concept, label in concepts:
        print(f"{concept} ({label})")

    question_type = input("\nEnter question type (one_word/true_false/3_mark/5_mark): ").strip()
    questions = generate_questions(summary, concepts, question_type)
    print("\n=== Generated Questions ===")
    print(questions)

if __name__ == "__main__":
    main()
